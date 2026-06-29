from __future__ import annotations

import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output"
ASSET_DIR = OUT_DIR / "ppt_assets"
PPT_PATH = OUT_DIR / "回声——复赛答辩PPT.pptx"

PRIMARY = RGBColor(0xE8, 0x94, 0x3A)
PRIMARY_DARK = RGBColor(0xD1, 0x7A, 0x2A)
PRIMARY_LIGHT = RGBColor(0xF0, 0xB8, 0x6C)
BG = RGBColor(0xF5, 0xF0, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x3C, 0x24, 0x15)
TEXT_LIGHT = RGBColor(0x6B, 0x5D, 0x53)
LINE = RGBColor(0xE8, 0xE3, 0xDD)
TAUPE = RGBColor(0xB5, 0xA9, 0x9A)
DARK_BG = RGBColor(0x42, 0x1D, 0x10)
MID_BROWN = RGBColor(0x6A, 0x41, 0x2B)

TITLE_FONT = "SimHei"
BODY_FONT = "Microsoft YaHei"
EN_FONT = "Arial"


def rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_glow(slide, left, top, width, height, color: RGBColor, transparency: float = 0.55) -> None:
    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
    glow.fill.solid()
    glow.fill.fore_color.rgb = color
    glow.fill.transparency = transparency
    glow.line.fill.background()


def add_soft_background(slide, dark: bool = False) -> None:
    set_slide_bg(slide, DARK_BG if dark else BG)
    if dark:
        add_glow(slide, Inches(-0.8), Inches(3.4), Inches(5.2), Inches(4.2), PRIMARY, 0.82)
        add_glow(slide, Inches(8.4), Inches(2.1), Inches(4.8), Inches(4.8), PRIMARY_LIGHT, 0.88)
    else:
        add_glow(slide, Inches(-1.0), Inches(6.18), Inches(3.55), Inches(1.52), PRIMARY_LIGHT, 0.94)
        add_glow(slide, Inches(9.75), Inches(6.16), Inches(3.35), Inches(1.50), PRIMARY, 0.95)


def add_wave_image(slide, image_path: Path, left, top, width, height) -> None:
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def add_title(slide, text: str, y: float = 0.36, size: int = 30, color: RGBColor = TEXT) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(7.2), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = TITLE_FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def add_subtitle(slide, text: str, y: float = 0.9, size: int = 14, color: RGBColor = TEXT_LIGHT) -> None:
    box = slide.shapes.add_textbox(Inches(0.82), Inches(y), Inches(11.0), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color


def add_line(slide, x1, y1, x2, y2, color: RGBColor = LINE, width_pt: float = 1.5, dash=None) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    if dash:
        line.line.dash_style = dash


def add_shadow_card(slide, left, top, width, height, fill_color: RGBColor = WHITE, radius=True,
                    line_color: RGBColor | None = None, shadow=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    if shadow:
        shadow_shape = slide.shapes.add_shape(
            shape_type,
            left + Inches(0.10),
            top + Inches(0.10),
            width,
            height,
        )
        shadow_shape.fill.solid()
        shadow_shape.fill.fore_color.rgb = RGBColor(0xA8, 0x8D, 0x73)
        shadow_shape.fill.transparency = 0.82
        shadow_shape.line.fill.background()
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    if line_color is None:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text: str, *,
                 font_name: str = BODY_FONT, size: int = 16, color: RGBColor = TEXT,
                 bold: bool = False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                 margin=(0.12, 0.12, 0.12, 0.12), line_spacing: float = 1.18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin[0])
    tf.margin_right = Inches(margin[1])
    tf.margin_top = Inches(margin[2])
    tf.margin_bottom = Inches(margin[3])
    paragraphs = text.split("\n")
    tf.clear()
    for idx, para_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = para_text
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_kicker(slide, left, top, width, label: str, size: int = 14, color: RGBColor = PRIMARY) -> None:
    box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.name = EN_FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def add_badge(slide, left, top, width, height, text: str, bg_color: RGBColor = PRIMARY, text_color: RGBColor = WHITE,
              font_size: int = 12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shape


def add_icon_circle(slide, left, top, diameter, icon: str, icon_size: int = 18, fill_color: RGBColor = PRIMARY) -> None:
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, diameter, diameter)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon
    run.font.name = BODY_FONT
    run.font.size = Pt(icon_size)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_placeholder_box(slide, left, top, width, height, label: str, aspect_label: str = "16:9") -> None:
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xFB, 0xF8, 0xF3)
    rect.line.color.rgb = RGBColor(0xBD, 0xB7, 0xB0)
    rect.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    rect.line.width = Pt(2)
    tf = rect.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{label}\n{aspect_label} 占位框"
    run.font.name = BODY_FONT
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT_LIGHT


def add_phone_placeholder(slide, left, top, width, height, label: str) -> None:
    phone = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    phone.fill.solid()
    phone.fill.fore_color.rgb = RGBColor(0xFD, 0xFB, 0xF8)
    phone.line.color.rgb = RGBColor(0xC9, 0xC0, 0xB6)
    phone.line.width = Pt(2)

    notch = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                   left + width * 0.34, top + Inches(0.08),
                                   width * 0.32, Inches(0.12))
    notch.fill.solid()
    notch.fill.fore_color.rgb = RGBColor(0xD9, 0xD2, 0xCA)
    notch.line.fill.background()

    screen = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                    left + Inches(0.14), top + Inches(0.26),
                                    width - Inches(0.28), height - Inches(0.48))
    screen.fill.solid()
    screen.fill.fore_color.rgb = BG
    screen.line.color.rgb = LINE
    screen.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    screen.line.width = Pt(1.6)
    tf = screen.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = BODY_FONT
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT_LIGHT


def add_annotation_list(slide, left, top, width, items: list[str], line_targets: list[tuple[float, float]]) -> None:
    y = top
    for idx, item in enumerate(items, start=1):
        num_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, y + Inches(0.02), Inches(0.28), Inches(0.28))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = PRIMARY
        num_box.line.fill.background()
        tf = num_box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(idx)
        run.font.name = EN_FONT
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE

        add_text_box(
            slide,
            left + Inches(0.36),
            y,
            width - Inches(0.36),
            Inches(0.52),
            item,
            size=14,
            color=TEXT,
            margin=(0, 0, 0, 0),
        )
        tx, ty = line_targets[idx - 1]
        add_line(
            slide,
            left - Inches(0.10),
            y + Inches(0.18),
            Inches(tx),
            Inches(ty),
            color=RGBColor(0xC9, 0xB1, 0x94),
            width_pt=1.2,
        )
        y += Inches(0.72)


def build_wave_image(path: Path, *, dark: bool = False, style: str = "cover") -> None:
    width, height = 1600, 900
    image = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(image)

    if style == "poster":
        center_x = width - 230
        center_y = 335
        for r in range(60, 240, 12):
            bbox = (center_x - r, center_y - r, center_x + r, center_y + r)
            color = (235, 180, 75, max(24, 175 - r // 2))
            draw.arc(bbox, start=18, end=342, fill=color, width=4)
        prev = None
        for x in range(40, center_x - 10, 8):
            progress = x / center_x
            amp = 7 + (progress ** 2.2) * 95
            freq = 0.045 + progress * 0.03
            y = center_y + math.sin(x * freq) * amp * (0.7 + progress * 0.5)
            if prev:
                draw.line([prev, (x, y)], fill=(244, 154, 58, 190), width=3)
            prev = (x, y)
        prev = None
        for x in range(40, center_x - 10, 8):
            progress = x / center_x
            amp = 3 + (progress ** 1.8) * 45
            freq = 0.058 + progress * 0.02
            y = center_y + math.cos(x * freq) * amp
            if prev:
                draw.line([prev, (x, y)], fill=(244, 197, 111, 170), width=2)
            prev = (x, y)
    else:
        base_y = height * (0.60 if style == "cover" else 0.52)
        color = (255, 255, 255, 120) if not dark else (238, 170, 77, 175)
        color2 = (255, 255, 255, 70) if not dark else (255, 205, 125, 108)
        for idx in range(13):
            pts = []
            for x in range(-40, width + 40, 12):
                progress = x / width
                amp = 42 + 28 * math.sin(progress * math.pi)
                y = base_y + math.sin(x * 0.012 + idx * 0.24) * amp + (idx - 6) * 10
                pts.append((x, y))
            draw.line(pts, fill=color if idx % 2 == 0 else color2, width=3)
        if dark:
            glow = Image.new("RGBA", (width, height), (0, 0, 0, 0))
            glow_draw = ImageDraw.Draw(glow)
            for idx in range(8):
                pts = []
                for x in range(-20, width + 20, 8):
                    y = base_y + math.sin(x * 0.016 + idx * 0.26) * (52 + idx * 5)
                    pts.append((x, y))
                glow_draw.line(pts, fill=(248, 184, 85, 28), width=10)
            glow = glow.filter(ImageFilter.GaussianBlur(8))
            image.alpha_composite(glow)

    image.save(path)


def build_cover_bg(path: Path) -> None:
    width, height = 1600, 900
    img = Image.new("RGB", (width, height), "#F0E7DA")
    px = img.load()
    for y in range(height):
        t = y / (height - 1)
        r = int(240 * (1 - t) + 176 * t)
        g = int(231 * (1 - t) + 98 * t)
        b = int(218 * (1 - t) + 36 * t)
        for x in range(width):
            px[x, y] = (r, g, b)
    glow = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(glow)
    draw.ellipse((820, 360, 1500, 980), fill=(255, 214, 164, 55))
    draw.ellipse((-240, 520, 720, 1180), fill=(255, 255, 255, 35))
    glow = glow.filter(ImageFilter.GaussianBlur(26))
    img = Image.alpha_composite(img.convert("RGBA"), glow)
    img.convert("RGB").save(path)


def build_dark_bg(path: Path) -> None:
    width, height = 1600, 900
    img = Image.new("RGB", (width, height), "#3A170D")
    px = img.load()
    for y in range(height):
        for x in range(width):
            dx = (x - width * 0.52) / width
            dy = (y - height * 0.52) / height
            glow = max(0.0, 1.0 - (dx * dx * 4.4 + dy * dy * 6.8))
            r = int(58 + glow * 40)
            g = int(23 + glow * 20)
            b = int(13 + glow * 8)
            px[x, y] = (r, g, b)
    img.save(path)


def ensure_assets() -> dict[str, Path]:
    OUT_DIR.mkdir(exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    paths = {
        "cover_bg": ASSET_DIR / "cover_bg.png",
        "cover_wave": ASSET_DIR / "cover_wave.png",
        "poster_wave": ASSET_DIR / "poster_wave.png",
        "final_bg": ASSET_DIR / "final_bg.png",
        "final_wave": ASSET_DIR / "final_wave.png",
    }
    build_cover_bg(paths["cover_bg"])
    build_wave_image(paths["cover_wave"], style="cover")
    build_wave_image(paths["poster_wave"], style="poster")
    build_dark_bg(paths["final_bg"])
    build_wave_image(paths["final_wave"], dark=True, style="final")
    return paths


def cover_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(assets["cover_bg"]), 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_wave_image(slide, assets["cover_wave"], Inches(0), Inches(3.42), Inches(13.33), Inches(2.1))
    add_text_box(slide, Inches(3.5), Inches(2.34), Inches(6.3), Inches(0.75), "回声", font_name=TITLE_FONT,
                 size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.45), Inches(3.28), Inches(6.5), Inches(0.45), "—— 留住即将消失的声音",
                 font_name=BODY_FONT, size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.05), Inches(5.82), Inches(5.2), Inches(0.35), "团队名称：不要操之过急",
                 font_name=BODY_FONT, size=17, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.4), Inches(6.28), Inches(6.7), Inches(0.35),
                 "2026中国高校计算机大赛 - AIGC创新赛 · 应用赛道",
                 font_name=EN_FONT, size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(10.55), Inches(6.92), Inches(2.3), Inches(0.18), "vivo BlueLM 端云协同",
                 font_name=EN_FONT, size=10, color=WHITE, align=PP_ALIGN.RIGHT)


def team_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "团队介绍")

    main = add_shadow_card(slide, Inches(0.8), Inches(1.55), Inches(5.45), Inches(4.85))
    add_text_box(slide, Inches(2.0), Inches(2.06), Inches(3.2), Inches(0.3), "团队名称：不要操之过急",
                 size=18, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, Inches(2.5), Inches(3.0), Inches(4.2), Inches(3.0), PRIMARY, 2.8)
    add_text_box(slide, Inches(2.0), Inches(3.78), Inches(3.2), Inches(0.6), "蓝海瑞",
                 font_name=TITLE_FONT, size=28, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.75), Inches(4.42), Inches(3.8), Inches(0.34), "华中科技大学 | 单人参赛",
                 size=16, color=TEXT, align=PP_ALIGN.CENTER)
    note = add_shadow_card(slide, Inches(1.15), Inches(5.52), Inches(4.75), Inches(0.5),
                           fill_color=RGBColor(0xFA, 0xF1, 0xE6), line_color=RGBColor(0xFA, 0xF1, 0xE6), shadow=False)
    add_text_box(slide, Inches(1.35), Inches(5.60), Inches(4.35), Inches(0.22),
                 "以上所有角色均由唯一成员独立承担",
                 size=14, color=PRIMARY_DARK, align=PP_ALIGN.CENTER, bold=True)

    role_cards = [
        ("产品策划", "需求分析、功能定义与整体规划"),
        ("UI / 交互设计", "界面视觉风格与用户体验流程设计"),
        ("前端开发", "网页与移动端页面的代码实现"),
        ("服务端开发", "后端逻辑、数据库架构与接口开发"),
    ]
    positions = [
        (Inches(6.8), Inches(1.75)),
        (Inches(10.0), Inches(1.75)),
        (Inches(6.8), Inches(4.0)),
        (Inches(10.0), Inches(4.0)),
    ]
    icons = ["策", "设", "前", "后"]
    for (title, desc), (left, top), icon in zip(role_cards, positions, icons):
        add_shadow_card(slide, left, top, Inches(2.85), Inches(1.95))
        add_icon_circle(slide, left + Inches(0.18), top + Inches(0.18), Inches(0.34), icon, 14)
        add_text_box(slide, left + Inches(0.62), top + Inches(0.20), Inches(1.9), Inches(0.3), title,
                     size=16, color=TEXT, bold=True)
        add_text_box(slide, left + Inches(0.24), top + Inches(0.74), Inches(2.3), Inches(0.78), desc,
                     size=13, color=TEXT_LIGHT)
    add_text_box(slide, Inches(0.92), Inches(6.48), Inches(11.4), Inches(0.3),
                 "一个人，三周时间，从零到一。因为有些声音，不该只留在记忆里。",
                 size=13, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)


def intro_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "作品简介")
    add_text_box(slide, Inches(1.05), Inches(1.45), Inches(11.2), Inches(1.05),
                 "「回声」是一款 AI 声音传承应用。用户只需与长辈聊一次天，即可沉淀一份可持续互动的声音档案；此后无论身处何地，都能继续听见熟悉的语气、口头禅与关怀。",
                 size=20, color=TEXT, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.24)
    cards = [
        ("一次聊天", "15-30 分钟自然对话采集\nAI 主动抛出话题\n不需要刻意朗读台词"),
        ("AI 建模", "Demo 版走上传 + 模拟训练\nAPK 规划接入 BlueLM 3B 端侧能力\n保留语调与说话习惯"),
        ("永续对话", "随时和“虚拟长辈”聊天\n支持文字与语音双入口\n自动沉淀为回忆时间线"),
    ]
    for idx, (title, desc) in enumerate(cards):
        left = Inches(0.82 + idx * 4.15)
        add_shadow_card(slide, left, Inches(3.1), Inches(3.55), Inches(2.35),
                        line_color=PRIMARY if idx != 1 else PRIMARY_LIGHT)
        add_badge(slide, left + Inches(0.18), Inches(3.26), Inches(0.72), Inches(0.28), f"0{idx + 1}")
        add_text_box(slide, left + Inches(0.18), Inches(3.62), Inches(3.0), Inches(0.36), title,
                     size=20, color=TEXT, bold=True)
        add_text_box(slide, left + Inches(0.18), Inches(4.08), Inches(3.0), Inches(1.05), desc,
                     size=14, color=TEXT_LIGHT)
    add_text_box(slide, Inches(1.6), Inches(6.35), Inches(10.0), Inches(0.28),
                 "核心创新：把语音 AI 从“效率工具”重新定义为“情感传承载体”",
                 size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


def quote_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, TAUPE)
    add_glow(slide, Inches(8.85), Inches(5.72), Inches(3.45), Inches(1.42), PRIMARY, 0.93)
    add_text_box(slide, Inches(2.1), Inches(1.8), Inches(9.2), Inches(0.55), "不是让机器学会说话，",
                 font_name=TITLE_FONT, size=30, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.0), Inches(2.48), Inches(9.3), Inches(0.55), "而是让爱可以续约。",
                 font_name=TITLE_FONT, size=30, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.2), Inches(3.45), Inches(6.8), Inches(0.35), "技术只是手段，留住声音里的温度才是目的。",
                 size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.45), Inches(5.15), Inches(10.4), Inches(0.85),
                 "“子欲养而亲不待，是这代人最深的遗憾。我们无法延长生命的长度，但可以让声音跨越时间的宽度。”",
                 size=14, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.2)


def why_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "为什么是「回声」？")
    data_cards = [
        ("130+", "种方言仍在流失", "方言不只是交流工具，更是地域文化与家庭记忆的活化石。"),
        ("3 亿+", "60 岁以上人口", "祖孙分隔两地已经成为常态，很多家庭没有足够时间系统留存声音。"),
        ("0", "款成熟的大众情感产品", "语音克隆不缺技术 demo，缺的是普通人真正愿意长期使用的产品形态。"),
    ]
    for idx, (num, title, desc) in enumerate(data_cards):
        left = Inches(0.82 + idx * 4.12)
        add_shadow_card(slide, left, Inches(1.55), Inches(3.58), Inches(2.78))
        add_text_box(slide, left + Inches(0.2), Inches(1.82), Inches(1.55), Inches(0.56), num,
                     font_name=EN_FONT, size=30, color=PRIMARY, bold=True)
        add_text_box(slide, left + Inches(1.55), Inches(1.92), Inches(1.5), Inches(0.4), title,
                     size=15, color=TEXT, bold=True)
        add_text_box(slide, left + Inches(0.2), Inches(2.62), Inches(3.05), Inches(1.1), desc,
                     size=14, color=TEXT_LIGHT)
    add_shadow_card(slide, Inches(0.82), Inches(5.05), Inches(11.7), Inches(1.22), fill_color=MID_BROWN,
                    line_color=MID_BROWN)
    add_text_box(slide, Inches(1.12), Inches(5.25), Inches(11.0), Inches(0.22), "核心判断：从“静态留存”走向“动态延续”",
                 size=18, color=PRIMARY_LIGHT, bold=True)
    add_text_box(slide, Inches(1.12), Inches(5.64), Inches(10.95), Inches(0.42),
                 "不是没有人想过留住声音，而是这件事过去太难、太贵、太不像一个能被日常使用的产品。大模型与端云协同，让它第一次有机会进入普通家庭。",
                 size=15, color=WHITE)


def design_values_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, TAUPE)
    add_text_box(slide, Inches(0.82), Inches(0.45), Inches(5.4), Inches(0.48), "设计理念，贯穿每一处细节",
                 font_name=TITLE_FONT, size=28, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.82), Inches(1.06), Inches(10.6), Inches(0.32),
                 "产品围绕三个核心词层层落地，让“让爱续约”不是口号，而是能被感知的体验。",
                 size=15, color=WHITE)
    cards = [
        ("温暖", "VISUAL · 视觉落地", "暖橙 + 米白 + 深棕的旧相册配色，让 AI 不再是冷色科技感，而更像家里的台灯与晚饭时分。"),
        ("简单", "INTERACTION · 交互落地", "摒弃“说明书式录音”，只保留“陪长辈聊一次天”的自然动作。录制、训练、进入对话全链路都尽量少解释。"),
        ("私密", "TECH · 技术落地", "原始录音保存在本地目录，Web 端云调用失败时自动降级到本地引擎；APK 环境再下沉到 BlueLM 3B 端侧桥接。"),
    ]
    for idx, (title, subtitle, desc) in enumerate(cards):
        left = Inches(0.82 + idx * 4.12)
        add_shadow_card(slide, left, Inches(2.05), Inches(3.55), Inches(4.05))
        add_icon_circle(slide, left + Inches(1.43), Inches(2.38), Inches(0.62), ["暖", "简", "密"][idx], 20)
        add_text_box(slide, left + Inches(0.4), Inches(3.28), Inches(2.8), Inches(0.38), title,
                     size=22, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.28), Inches(3.78), Inches(3.0), Inches(0.25), subtitle,
                     font_name=EN_FONT, size=14, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.28), Inches(4.40), Inches(3.0), Inches(1.14), desc,
                     size=14, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.12), Inches(6.52), Inches(7.15), Inches(0.24),
                 "好的设计是隐形的。用户感受到的是温暖，而不是技术术语。",
                 size=14, color=WHITE, align=PP_ALIGN.CENTER)


def core_features_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "核心功能")
    top_card = add_shadow_card(slide, Inches(0.82), Inches(1.46), Inches(11.72), Inches(1.08))
    add_text_box(slide, Inches(1.08), Inches(1.72), Inches(10.95), Inches(0.52),
                 "通过 AI 引导式自然对话完成声音采集，结合云端大模型回复、本地语义兜底与端侧 BlueLM 3B 规划路径，形成完整的“采集 → 建模 → 陪伴”闭环。",
                 size=17, color=TEXT)
    items = [
        ("01 / 一键录制", "输入称呼 → 10 条 AI 引导话题 → 最长 10 分钟录音 → 实时声波反馈 → 不要求用户朗读固定脚本。"),
        ("02 / AI 建模", "上传后进入 processing / training / ready 状态推进。Demo 版约 25 秒完成模拟训练进度，验证完整产品链路。"),
        ("03 / 永续对话", "对话页支持文字输入 + 按住说话；端侧可用时优先走 BlueLM 3B，失败后静默降级到云端与本地引擎。"),
    ]
    for idx, (title, desc) in enumerate(items):
        left = Inches(0.82 + idx * 4.15)
        add_shadow_card(slide, left, Inches(3.18), Inches(3.55), Inches(2.96))
        add_badge(slide, left + Inches(0.22), Inches(3.48), Inches(0.8), Inches(0.30), ["录", "模", "聊"][idx])
        add_text_box(slide, left + Inches(0.22), Inches(3.84), Inches(3.0), Inches(0.35), title,
                     size=19, color=TEXT, bold=True)
        add_text_box(slide, left + Inches(0.20), Inches(4.58), Inches(3.08), Inches(1.42), desc,
                     size=13, color=TEXT_LIGHT, margin=(0, 0, 0, 0))
    add_text_box(slide, Inches(3.4), Inches(6.46), Inches(6.6), Inches(0.22),
                 "全流程约 1 分钟可体验到“从录音到开聊”的第一轮闭环。",
                 size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


def screenshot_slide(prs: Presentation, title: str, placeholder: str, notes: list[str], footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, title)
    add_phone_placeholder(slide, Inches(1.05), Inches(1.46), Inches(4.2), Inches(4.9), placeholder)
    add_annotation_list(
        slide,
        Inches(6.0),
        Inches(1.75),
        Inches(5.65),
        notes,
        [(5.0, 2.1), (4.55, 3.0), (4.9, 4.0), (4.7, 5.1)],
    )
    add_text_box(slide, Inches(1.1), Inches(6.48), Inches(11.0), Inches(0.22), footer,
                 size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


def flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "用户交互流程")
    add_shadow_card(slide, Inches(0.82), Inches(1.55), Inches(11.75), Inches(1.4), fill_color=RGBColor(0xFF, 0xF4, 0xE8),
                    line_color=RGBColor(0xFF, 0xF4, 0xE8))
    add_text_box(slide, Inches(1.02), Inches(1.78), Inches(2.25), Inches(0.22), "流程一 · 声音采集",
                 size=16, color=PRIMARY, bold=True)
    add_text_box(slide, Inches(2.82), Inches(2.12), Inches(9.4), Inches(0.46),
                 "🏠 打开应用 → ➕ 点击“+” → ✏️ 输入称呼 → 💬 AI 抛出引导话题 → 🎙️ 自然录制 → ⚙️ 上传并轮询训练状态 → ✅ 就绪通知",
                 size=15, color=TEXT)

    add_shadow_card(slide, Inches(0.82), Inches(3.32), Inches(11.75), Inches(1.42), fill_color=MID_BROWN,
                    line_color=MID_BROWN)
    add_text_box(slide, Inches(1.02), Inches(3.55), Inches(2.25), Inches(0.22), "流程二 · 永续对话",
                 size=16, color=PRIMARY_LIGHT, bold=True)
    add_text_box(slide, Inches(2.82), Inches(3.90), Inches(9.2), Inches(0.46),
                 "🏠 首页点击声音卡片 → 💬 文字或语音输入 → 🧠 端云调度决策路径 → 🤖 生成拟人化回复 → 🔊 Edge TTS / 浏览器朗读 → 💾 自动写入回忆时间线",
                 size=15, color=WHITE)

    mini = [
        ("对话引导式采集", "通过聊天而不是念稿，降低长辈面对麦克风的压力。"),
        ("端云多层降级", "端侧不可用时自动走云端，云端失败时回落本地引擎。"),
        ("话题推荐区", "首次对话展示常用话题 chip，点一下就能发出第一句。"),
        ("回忆自动沉淀", "每次对话完成后保存到本地 SQLite 与时间线页面。"),
    ]
    for idx, (title, desc) in enumerate(mini):
        left = Inches(0.82 + idx * 3.0)
        add_shadow_card(slide, left, Inches(5.18), Inches(2.72), Inches(1.34),
                        fill_color=RGBColor(0xFB, 0xF4, 0xEA), line_color=RGBColor(0xFB, 0xF4, 0xEA), shadow=False)
        add_text_box(slide, left + Inches(0.16), Inches(5.38), Inches(2.4), Inches(0.25), title,
                     size=15, color=TEXT, bold=True)
        add_text_box(slide, left + Inches(0.16), Inches(5.72), Inches(2.42), Inches(0.55), desc,
                     size=12, color=TEXT_LIGHT)
    add_text_box(slide, Inches(3.2), Inches(6.65), Inches(7.0), Inches(0.22),
                 "从打开应用到发出第一条消息，全程不超过 60 秒。",
                 size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


def llm_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "大模型的具体应用与端云调度 ★", color=PRIMARY)
    add_subtitle(slide, "本章节对应赛事“大模型应用能力”评分项，内容已按当前代码仓库的真实实现校准。")

    card_specs = [
        (
            "核心对话大模型 —— vivo 蓝心平台 API",
            [
                "调用位置：app.py → call_bluelm_cloud()",
                "System Prompt：65-80 岁老人角色、20-80 字口语化回复、带关怀语，不暴露 AI 身份。",
                "多模型轮询：Volc-DeepSeek-V3.2 → qwen3.5-plus → Doubao-Seed-2.0-mini，单次超时 10 秒。",
                "兜底保障：云端失败即返回 None，由 generate_mock_reply() 接管，确保演示稳定。",
            ],
        ),
        (
            "语音合成 —— Edge TTS + 浏览器 Web Speech API",
            [
                "调用位置：app.py → generate_tts_audio() / chat.js → _speakWithBrowserTTS()",
                "音色自动匹配：称呼含“奶/婆/妈”默认女声，含“爷/公/爸”默认男声。",
                "参数配置：女声 -15% 语速、-8Hz 音调；男声 -20% 语速、-12Hz 音调，模拟长辈节奏。",
                "离线兜底：Edge TTS 不可用时自动回落系统自带 Web Speech API。",
            ],
        ),
        (
            "端侧声纹与推理路径 —— BlueLM 3B 决赛规划",
            [
                "技术路线：android-bridge/WebAppInterface.java 已定义 initModel() / chat() / isNative() 完整桥接。",
                "调度策略：chat.js 中 _callOnDevice() 优先尝试端侧 3B，5 秒无响应即熔断降级到云端。",
                "隐私主张：Web 公开版验证业务逻辑，APK 版本负责把真正敏感推理下沉到手机侧。",
                "Demo 说明：当前复赛阶段以“云端 + 本地兜底”完整跑通闭环，端侧能力在 APK 工程中预留。 ",
            ],
        ),
    ]

    top = 1.35
    for title, bullets in card_specs:
        add_shadow_card(slide, Inches(0.82), Inches(top), Inches(11.72), Inches(1.32))
        add_text_box(slide, Inches(1.0), Inches(top + 0.14), Inches(11.2), Inches(0.22), title,
                     size=16, color=TEXT, bold=True)
        for idx, bullet in enumerate(bullets):
            add_text_box(slide, Inches(1.0), Inches(top + 0.44 + idx * 0.19), Inches(10.95), Inches(0.18),
                         f"• {bullet}", size=11, color=TEXT_LIGHT)
        top += 1.52

    add_placeholder_box(slide, Inches(0.82), Inches(6.03), Inches(5.3), Inches(0.95), "05-代码.png", "代码截图")
    add_text_box(slide, Inches(6.35), Inches(6.05), Inches(5.95), Inches(0.74),
                 "建议最终替换为 app.py 中 call_bluelm_cloud() 的 VS Code 截图，包含多模型轮询与 System Prompt 构造逻辑。\n这样评委在 PPT 与代码包之间能形成一一对应。",
                 size=12, color=TEXT_LIGHT)
    add_text_box(slide, Inches(0.95), Inches(6.73), Inches(11.2), Inches(0.18),
                 "* 复赛 Demo 当前采用“云端 API + 本地兜底”验证核心逻辑；决赛阶段将继续强化端侧能力。",
                 size=9, color=TEXT_LIGHT, align=PP_ALIGN.LEFT)


def innovation_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "五大创新，重新定义“语音 AI”")
    cards = [
        ("01", "方案创新", "把语音 AI 从“生产力工具”转向“情感传承载体”，关注的是陪伴，不只是效率。"),
        ("02", "交互创新", "用引导式聊天完成采集，避免固定朗读文本带来的表演感与使用门槛。"),
        ("03", "功能创新", "从静态录音升级到动态永续对话，让声音真正成为可继续交流的数字人格。"),
        ("04", "文化创新", "把方言保护从专业采样拉回日常家庭场景，让每个人都能成为文化保存者。"),
        ("05", "隐私创新", "坚持端侧优先、云端可降级、本地永兜底的架构，用户始终掌握数据主权。"),
    ]
    positions = [
        (Inches(0.82), Inches(1.7), Inches(3.55), Inches(2.05)),
        (Inches(4.52), Inches(1.7), Inches(3.55), Inches(2.05)),
        (Inches(8.22), Inches(1.7), Inches(4.3), Inches(2.05)),
        (Inches(0.82), Inches(4.05), Inches(5.95), Inches(2.1)),
        (Inches(6.95), Inches(4.05), Inches(5.57), Inches(2.1)),
    ]
    for (num, title, desc), (left, top, width, height) in zip(cards, positions):
        add_shadow_card(slide, left, top, width, height)
        add_text_box(slide, left + Inches(0.18), top + Inches(0.18), Inches(1.0), Inches(0.5), num,
                     font_name=EN_FONT, size=30, color=PRIMARY, bold=True)
        add_text_box(slide, left + Inches(1.16), top + Inches(0.26), width - Inches(1.4), Inches(0.3), title,
                     size=18, color=TEXT, bold=True)
        add_text_box(slide, left + Inches(1.16), top + Inches(0.74), width - Inches(1.5), height - Inches(0.92), desc,
                     size=14, color=TEXT_LIGHT)


def persona_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "为谁而做？用在哪里？")

    add_shadow_card(slide, Inches(0.82), Inches(1.55), Inches(6.1), Inches(4.65))
    add_text_box(slide, Inches(1.02), Inches(1.78), Inches(2.6), Inches(0.24), "目标用户画像",
                 size=18, color=TEXT, bold=True)
    add_text_box(slide, Inches(1.02), Inches(2.16), Inches(5.5), Inches(0.44),
                 "核心用户：18-35 岁异地求学 / 工作人群，与祖辈、父母辈感情深厚，愿意为“情感留存”付出时间。",
                 size=15, color=TEXT)
    pain_points = [
        "害怕遗忘：“我好怕有一天，真的记不清她的声音。”",
        "不知聊什么：电话总在“吃了没、冷不冷”里很快结束。",
        "子欲养而亲不待：想陪伴、想记录，但不知道从哪一步开始。",
    ]
    y = 2.78
    for item in pain_points:
        add_text_box(slide, Inches(1.12), Inches(y), Inches(5.4), Inches(0.36), f"• {item}",
                     size=14, color=TEXT_LIGHT)
        y += 0.58
    add_shadow_card(slide, Inches(1.02), Inches(4.78), Inches(5.65), Inches(0.98),
                    fill_color=RGBColor(0xFB, 0xF4, 0xEA), line_color=RGBColor(0xFB, 0xF4, 0xEA), shadow=False)
    add_text_box(slide, Inches(1.22), Inches(5.00), Inches(5.15), Inches(0.48),
                 "用户不是在说“我需要一个工具”，而是在说“我需要一种方式，让思念还能被触碰”。",
                 size=15, color=PRIMARY_DARK, align=PP_ALIGN.CENTER, bold=True)

    add_shadow_card(slide, Inches(7.12), Inches(1.55), Inches(5.4), Inches(4.65))
    add_text_box(slide, Inches(7.32), Inches(1.78), Inches(2.8), Inches(0.24), "典型使用场景",
                 size=18, color=TEXT, bold=True)
    scenes = [
        "🌙 异乡深夜：加班到很晚，点开回声，听到熟悉的那句“别太累了”。",
        "🧧 逢年过节：把新年祝福录下来，让它不只在当年存在。",
        "🍳 学一道菜：把“怎么做那道菜”从视频教程变成真正熟悉的口吻。",
        "👶 给下一代：把家族里最温柔的声音，留给未来还没来得及见面的人。",
    ]
    y = 2.22
    for item in scenes:
        add_text_box(slide, Inches(7.36), Inches(y), Inches(4.8), Inches(0.58), item,
                     size=14, color=TEXT_LIGHT)
        y += 0.85
    add_text_box(slide, Inches(1.7), Inches(6.24), Inches(9.9), Inches(0.24),
                 "回声不是生产力工具，更像一份面向失去的情感保险。",
                 size=15, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)


def competition_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "为什么是回声？")
    add_shadow_card(slide, Inches(0.82), Inches(1.42), Inches(11.72), Inches(2.72))
    headers = ["维度", "录音 APP", "语音助手", "AI 语音克隆工具", "回声"]
    col_x = [0.98, 2.25, 4.18, 6.35, 9.45]
    col_w = [1.15, 1.75, 1.8, 2.75, 2.45]
    for x, w, head in zip(col_x, col_w, headers):
        add_text_box(slide, Inches(x), Inches(1.64), Inches(w), Inches(0.28), head,
                     size=15, color=PRIMARY if head == "回声" else TEXT, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=(0, 0, 0, 0))
    rows = [
        ("核心价值", "单向回放", "指令执行", "技术 demo", "情感传承"),
        ("声音个性", "原声", "通用 TTS", "机械克隆", "人格化 AI"),
        ("交互方式", "播放 / 暂停", "命令式语音", "文本编辑", "自然对话"),
        ("情感连接", "弱", "弱", "弱", "强"),
        ("使用门槛", "低", "中", "高", "极低"),
    ]
    y = 1.98
    row_h = 0.43
    for row in rows:
        for x, w, cell in zip(col_x, col_w, row):
            add_text_box(slide, Inches(x), Inches(y), Inches(w), Inches(0.32), cell,
                         size=13, color=TEXT if cell != "情感传承" and cell != "人格化 AI" and cell != "自然对话" and cell != "强" and cell != "极低" else PRIMARY_DARK,
                         bold=cell in {"情感传承", "人格化 AI", "自然对话", "强", "极低"},
                         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=(0, 0, 0, 0))
        add_line(slide, Inches(1.0), Inches(y + row_h), Inches(12.0), Inches(y + row_h), LINE, 0.8)
        y += row_h

    advantages = [
        ("品类创新", "不是“更好的录音”，而是“另一种延续”。"),
        ("情感粘性", "留存来自关系，而不是单次工具效率。"),
        ("vivo 生态", "和 BlueLM 端云协同方向天然契合。"),
        ("隐私安全", "本地存储 + 多层降级，强调数据主权。"),
        ("传播自驱", "“听听奶奶的声音”天然自带分享冲动。"),
    ]
    for idx, (title, desc) in enumerate(advantages):
        left = Inches(0.82 + idx * 2.35)
        add_shadow_card(slide, left, Inches(4.42), Inches(2.1), Inches(1.72))
        add_text_box(slide, left + Inches(0.14), Inches(4.66), Inches(1.82), Inches(0.24), title,
                     size=14, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.12), Inches(5.02), Inches(1.86), Inches(0.62), desc,
                     size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


def architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_title(slide, "端云协同 · 多层降级 · 永不崩溃")
    layers = [
        ("第 1 层：云端 · vivo 蓝心平台 API", "处理复杂情感对话、多模型轮询与 System Prompt 驱动的拟人化回复。", PRIMARY),
        ("第 2 层：端侧 · BlueLM 3B SDK（APK 路径）", "Android WebView 已预留 JSBridge 调用 initModel() / chat()；当前 Demo 返回降级信号，真机 NPU 仍需联调验收。", MID_BROWN),
        ("第 3 层：本地 · Flask 降级引擎", "generate_mock_reply() 提供 10 类语义匹配回复；TTS 失败时再降级到浏览器 Web Speech API。", RGBColor(0xEA, 0xE2, 0xD6)),
    ]
    y = 1.55
    for idx, (title, desc, color) in enumerate(layers):
        add_shadow_card(slide, Inches(0.82), Inches(y), Inches(11.72), Inches(0.92), fill_color=color, line_color=color, shadow=False)
        add_text_box(slide, Inches(1.02), Inches(y + 0.16), Inches(4.2), Inches(0.22), title,
                     size=17, color=WHITE if idx < 2 else TEXT, bold=True)
        add_text_box(slide, Inches(5.15), Inches(y + 0.15), Inches(6.1), Inches(0.38), desc,
                     size=13, color=WHITE if idx < 2 else TEXT_LIGHT)
        y += 1.02

    add_shadow_card(slide, Inches(0.82), Inches(4.84), Inches(11.72), Inches(1.3))
    tech_lines = [
        "前端：原生 HTML5 / CSS3 / JavaScript，5 个核心页面，移动端 480px 优先。",
        "后端：Python Flask + SQLite WAL + threading.Lock，保障本地持久化与并发安全。",
        "AI 管线：call_bluelm_cloud() → generate_mock_reply() → _callOnDevice() 的端云协同调度。",
        "部署：本地 Flask 可运行，Hugging Face 仅作备用演示；Sealos 当前仍在部署准备中。",
    ]
    for idx, line in enumerate(tech_lines):
        add_text_box(slide, Inches(1.06), Inches(5.06 + idx * 0.22), Inches(10.95), Inches(0.18),
                     f"• {line}", size=13, color=TEXT_LIGHT)
    add_text_box(slide, Inches(2.0), Inches(6.24), Inches(9.2), Inches(0.24),
                 "当前复赛版以云端调用和本地兜底跑通闭环，端侧能力保留到 APK 路径继续联调。",
                 size=15, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)


def closing_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(assets["final_bg"]), 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_wave_image(slide, assets["final_wave"], Inches(0.3), Inches(2.0), Inches(12.7), Inches(2.7))
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11.8), Inches(0.78),
                 "有些声音，不该只留在记忆里。",
                 font_name=TITLE_FONT, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.65), Inches(4.35), Inches(2.0), Inches(0.35), "—— 回声",
                 font_name=BODY_FONT, size=20, color=PRIMARY_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.72), Inches(5.4), Inches(5.9), Inches(0.28),
                 "感谢各位评委老师的时间和宝贵意见",
                 size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.2), Inches(5.86), Inches(6.9), Inches(0.22),
                 "联系方式：blehsrtkk@163.com",
                 font_name=EN_FONT, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.85), Inches(6.18), Inches(9.6), Inches(0.22),
                 "在线体验：https://break66-echo-app.hf.space",
                 font_name=EN_FONT, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.1), Inches(6.86), Inches(9.1), Inches(0.18),
                 "2026 中国高校计算机大赛 - AIGC 创新赛 · 应用赛道 · 作品提交",
                 font_name=EN_FONT, size=10, color=WHITE, align=PP_ALIGN.CENTER)


def build_presentation() -> Path:
    assets = ensure_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover_slide(prs, assets)
    team_slide(prs)
    intro_slide(prs)
    quote_slide(prs)
    why_slide(prs)
    design_values_slide(prs)
    core_features_slide(prs)

    screenshot_slide(
        prs,
        "产品展示 · 首页",
        "01-首页.png",
        [
            "动态问候语根据时段切换：早上好 / 下午好 / 晚上好。",
            "声音档案卡片展示头像、名字与当前模型状态。",
            "悬浮“+”按钮在滚动时会降低透明度，停止后恢复。",
            "底部导航固定为 首页 / 发现 / 回忆 三入口。",
        ],
        "移动端 480px 优先设计，Web 公开版与 APK 版本保持一致的核心信息架构。",
    )
    screenshot_slide(
        prs,
        "产品展示 · 声音录制",
        "02-录制.png",
        [
            "称呼输入支持任意中文称谓，并自动复用上次使用记录。",
            "录音页基于 MediaRecorder + Web Audio API，实时显示声波与计时。",
            "AI 话题卡内置 10 条引导问题，可点击箭头切换下一题。",
            "录制状态机为 nameInput → recording → confirm → uploading → training → ready。",
        ],
        "最长录制 10 分钟；上传后通过 /api/voices/<id>/model-status 轮询训练进度。",
    )
    screenshot_slide(
        prs,
        "产品展示 · AI 对话",
        "03-对话.png",
        [
            "长辈回复以固定角色设定生成，保持语气一致与口语化表达。",
            "用户既可键入文字，也可按住说话发送语音消息。",
            "音频优先播放 Edge TTS；若失败，则自动回落浏览器 Web Speech API。",
            "首次进入页面展示话题 chip，帮助用户更快开始第一句。",
        ],
        "APK Demo 会先进入端侧调度入口，当前桥接返回降级信号后继续走云端 API 与本地引擎。",
    )
    screenshot_slide(
        prs,
        "产品展示 · 时光留声机",
        "04-回忆.png",
        [
            "统计卡片展示声音档案数、总对话次数与最近活跃情况。",
            "回忆时间线按时间倒序展示每次对话摘要与播放入口。",
            "话题标签云自动聚合“天气 / 家庭 / 美食 / 回忆”等高频主题。",
            "所有记录默认写入本地 SQLite，优先保证家庭数据可控。",
        ],
        "数据存储结构为 voices + conversations 两张表，便于本地检索、删除与迁移。",
    )
    flow_slide(prs)
    llm_slide(prs)
    innovation_slide(prs)
    persona_slide(prs)
    competition_slide(prs)
    architecture_slide(prs)
    closing_slide(prs, assets)

    prs.save(PPT_PATH)
    return PPT_PATH


if __name__ == "__main__":
    out = build_presentation()
    print(out)
